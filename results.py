import glob
import os
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from rdkit import Chem, DataStructs
from rdkit.Chem import Descriptors, AllChem, rdmolfiles
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

# Get the directory of the current script
script_dir = os.path.dirname(os.path.realpath(__file__))

# Change the current directory to the script directory
os.chdir(script_dir)

# Get a list of all output files in the Results subdirectory
output_files = glob.glob("Results/*.pdbqt")

# Create a list of tuples, each containing the file name, its best binding affinity, and its molecular weight
data = []
for file_name in output_files:
    ligand_id = file_name.replace('Results/', '').replace('_output.pdbqt', '')
    if ligand_id.lower().startswith('ligand'):
        continue
    with open(file_name, 'r') as file:
        for line in file:
            if line.startswith("REMARK VINA RESULT"):
                # The binding affinity is the second value on the line
                affinity = float(line.split()[3])
                # Read the molecule from a SDF file
                mol = Chem.SDMolSupplier(f'{script_dir}/Ligands/PDBQT/{ligand_id}.sdf')[0]
                # Calculate the molecular weight
                mol_weight = Descriptors.MolWt(mol)
                data.append((ligand_id, mol_weight, affinity, mol))
                # Stop after the first REMARK VINA RESULT line, as this is the best binding affinity
                break

# Create a DataFrame from the data
df = pd.DataFrame(data, columns=['LigandID', 'MolecularWeight', 'Affinity', 'Molecule'])

# Extract just the affinities from the DataFrame
affinities = df['Affinity']

# Calculate the mean and standard deviation
mean = np.mean(affinities)
std_dev = np.std(affinities)

# Exclude any affinities that are more than 2 standard deviations away from the mean
affinities = [x for x in affinities if (x > mean - 2*std_dev) and (x < mean + 2*std_dev)]

# Create a histogram of the binding affinities
plt.hist(affinities, bins=20, edgecolor='black')

# Add labels and title
plt.xlabel('Binding Affinities (kcal/mol)')
plt.ylabel('Frequency')
plt.title('Histogram of Binding Affinities')

# Display the plot
plt.savefig('histogram.png')

df = df.sort_values('Affinity').reset_index(drop=True)

# Print the DataFrame
print(df[:11])

df = df[df['Affinity']  < 0]


# Create a scatter plot of molecular weight vs. affinity for all rows
plt.scatter(df['MolecularWeight'], df['Affinity'], color='blue')

# Create a scatter plot of molecular weight vs. affinity for the top 10 rows
plt.scatter(df['MolecularWeight'][:10], df['Affinity'][:10], color='yellow')

plt.ylim(top=0)
# Add labels and title
plt.xlabel('Molecular Weight')
plt.ylabel('Affinity')
plt.title('Molecular Weight vs. Affinity')

# Display the plot
plt.savefig('scatter_plot.png')



##################


# Sort the DataFrame by the 'Affinity' column
df = df.sort_values(by='Affinity')

# Calculate the fingerprints for all molecules
df['Fingerprint'] = df['Molecule'].apply(lambda x: AllChem.GetMorganFingerprintAsBitVect(x, 2))

# Calculate the similarity between all pairs of fingerprints
similarity_matrix = []
for i in range(len(df)):
    similarities = []
    for j in range(len(df)):
        similarity = DataStructs.FingerprintSimilarity(df['Fingerprint'][i], df['Fingerprint'][j])
        similarities.append(similarity)
    similarity_matrix.append(similarities)

# Convert the similarity matrix to a DataFrame
df_similarity = pd.DataFrame(similarity_matrix, index=df['LigandID'], columns=df['LigandID'])

# Create a heatmap from the similarity DataFrame
plt.figure(figsize=(10, 8))
sns.heatmap(df_similarity, cmap='viridis')

# Save the plot
plt.savefig('heatmap.png')

# Get the top 10 molecules
top_10_molecules = df.nsmallest(10, 'Affinity')

# Save each molecule as a PDB file
for i, row in top_10_molecules.iterrows():
    # Use the LigandID from the dataframe
    ligand_id = row['LigandID']
    
    # Save the molecule as a PDB file with the LigandID as the filename
    rdmolfiles.MolToPDBFile(row['Molecule'], f"{ligand_id}.pdb")

# Print the SMILES string for each of the top 10 molecules
for i, row in top_10_molecules.iterrows():
    smiles = Chem.MolToSmiles(row['Molecule'])
    print(f"Molecule {i}: {smiles}")


# Create a PowerPoint presentation object
prs = Presentation()

# Define the width and height of the image
width = height = Inches(5.0)

# Drop the 'Molecule' column
df = df.drop(columns=['Molecule'])

# Convert the DataFrame to a list of lists
df_list = df.values.tolist()

# Add column names as the first row in df_list
df_list.insert(0, df.columns.tolist())

# Add a slide with a title and content layout
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Results"

# Define table data
rows, cols = len(df_list), len(df_list[0])
left = top = Inches(2.0)
width = height = Inches(6.0)

# Add table to slide
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Fill table cells with data
for r in range(rows):
    for c in range(cols):
        table.cell(r, c).text = str(df_list[r][c])


# Add the histogram image to the slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture('histogram.png', Inches(2), Inches(2), width, height)

# Add the scatter plot image to the slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture('scatter_plot.png', Inches(2), Inches(2), width, height)

# Add the heatmap image to the slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture('heatmap.png', Inches(2), Inches(2), width, height)

# Save the presentation
prs.save('results.pptx')